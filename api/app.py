import base64
from datetime import datetime
from fastapi import FastAPI, HTTPException, Request
from pydantic import BaseModel
import requests
import os
from pptx import Presentation
import logging
from typing import Optional, List
import uuid
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from starlette.middleware.sessions import SessionMiddleware
import uvicorn
from fastapi import File, UploadFile
import fitz
import torch
import easyocr
from PIL import Image
from docx import Document
import openpyxl



# Configuración básica de logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuración desde variables de entorno
OLLAMA_HOST = os.getenv("OLLAMA_HOST", "http://localhost:11434")
MODEL_NAME = os.getenv("MODEL_NAME", "gemma3:12b") # mistral, phi4, llama3, llama3.2, gemma3:12b, gemma3:27b

app = FastAPI(
    title="Chat Inteligente sin Búsqueda",
    description="API que genera respuestas utilizando un modelo de lenguaje sin búsqueda web.",
    version="1.0.0"
)

# Configuración de archivos estáticos y templates
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")

# Configuración del middleware de sesiones
app.add_middleware(SessionMiddleware, secret_key="vivan_las_practicas")

# Diccionario para almacenar los historiales de las sesiones (en memoria)
session_histories = {}

@app.get("/", include_in_schema=False)
async def chat_interface(request: Request):
    """
    Endpoint principal que sirve la interfaz web de chat con manejo de sesión
    
    Args:
        request (Request): Objeto de solicitud FastAPI
    
    Returns:
        TemplateResponse: Renderiza la plantilla index.html
    """
    # Recuperar o generar un session_id
    session_id = request.cookies.get("session_id")
    
    if not session_id:
        session_id = str(uuid.uuid4())  # Generar un ID de sesión nuevo si no existe

    # Limpiar historial anterior (si existiera)
    session_histories.pop(session_id, None)

    # Establecer la cookie de la sesión
    response = templates.TemplateResponse("index.html", {"request": request, "session_id": session_id})
    response.set_cookie("session_id", session_id)
    
    return response

class PromptRequest(BaseModel):
    """Modelo Pydantic para las solicitudes de generación de texto"""
    prompt: str
    max_tokens: Optional[int] = None
    file_type: Optional[str] = None
    file_name: Optional[str] = None
    file_text: Optional[str] = None  

class SimplifiedResponse(BaseModel):
    """Modelo Pydantic para las respuestas simplificadas"""
    model: str
    response: str
    sources: Optional[List[str]] = None
    context_used: bool = False


def read_pdf(file_path: str) -> str:
    """Lee el contenido de un archivo PDF y extrae el texto"""
    try:
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text("text")
        return text
    except Exception as e:
        logger.error(f"Error al leer el PDF: {str(e)}")
        return ""

def read_docx(file_path: str) -> str:
    """Lee el contenido de un archivo DOCX y extrae el texto"""
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error al leer el DOCX: {str(e)}")
        return ""
    
def read_txt(file_path: str) -> str:
    """Lee el contenido de un archivo TXT y extrae el texto"""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return text
    except Exception as e:
        logger.error(f"Error al leer el TXT: {str(e)}")
        return ""

def read_pptx(file_path: str) -> str:
    """Lee el contenido de un archivo PPTX y extrae el texto"""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Error al leer el PPTX: {str(e)}")
        return ""
    
def read_xlsx(file_path: str) -> str:
    """Lee el contenido de un archivo XLSX y extrae el texto de todas las celdas"""
    try:
        # Cargar el archivo Excel
        wb = openpyxl.load_workbook(file_path, data_only=True)
        sheet = wb.active  # Obtener la hoja activa

        text = ""
        # Iterar sobre todas las filas y columnas de la hoja activa
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value:
                    text += str(cell.value) + "\t"  # Usamos tabulador para separar los valores
            text += "\n"  # Nueva línea para cada fila

        return text.strip()  # Eliminar espacios adicionales al final

    except Exception as e:
        logger.error(f"Error al leer el XLSX: {str(e)}")
        return ""

def read_image(file_path: str) -> str:
    """Lee el contenido de una imagen JPG y extrae el texto usando OCR (easyOCR)."""
    try:
        # Crear un objeto Reader de easyOCR para el idioma deseado (puedes usar 'es' para español, por ejemplo)
        reader = easyocr.Reader(['es','en'])  # Usa ['es'] para español, ['en'] para inglés, o ambos.

        # Usar easyOCR para extraer el texto de la imagen
        result = reader.readtext(file_path)

        # Extraer el texto de la salida de easyOCR
        text = ""
        for detection in result:
            text += detection[1] + "\n"  # El texto extraído está en la segunda posición de la tupla

        return text.strip()  # Eliminar espacios adicionales al final

    except Exception as e:
        logger.error(f"Error al leer la imagen: {str(e)}")
        return ""

@app.post("/upload_file")
async def upload_file(file: UploadFile = File(...)):
    """
    Endpoint para cargar documentos (PDF, DOCX, TXT, PPTX) o imágenes (JPG, PNG, JPEG, WEBP).
    
    - Extrae texto de documentos.
    - Convierte imágenes a Base64.
    
    Args:
        file (UploadFile): Archivo a procesar.
    
    Returns:
        dict: Texto extraído o imagen en Base64.
    """
    try:
        file_ext = file.filename.split(".")[-1].lower()
        temp_file_path = f"temp_{uuid.uuid4()}.{file_ext}"

        with open(temp_file_path, "wb") as f:
            f.write(await file.read())
        
        match file_ext:
            case "pdf":
                # 📜 Leer texto de PDF
                text = read_pdf(temp_file_path)
            case "docx":
                # 📜 Leer texto de DOCX
                text = read_docx(temp_file_path)
            case "txt":
                # 📜 Leer texto de TXT
                text = read_txt(temp_file_path)
            case "pptx":
                # 📜 Leer texto de PPTX
                text = read_pptx(temp_file_path)
            case "xlsx":
                # 📜 Leer texto de XLSX
                text = read_xlsx(temp_file_path)
            case "jpg" | "png" | "jpeg" | "webp":
                # 🖼️ Leer texto de imagen
                text = read_image(temp_file_path)
            case _:
                os.remove(temp_file_path)
                raise HTTPException(status_code=400, detail="Formato de archivo no soportado")

        # Eliminar archivo temporal después de procesarlo
        os.remove(temp_file_path)

        # Devolver respuesta con el texto extraído y el nombre del archivo
        return {"file_type": "documento" if file_ext in ["pdf", "docx", "txt", "pptx", "xlsx"] else "imagen", "file_text": text[:8000], "file_name": file.filename}

        
    except Exception as e:
        logger.error(f"Error al procesar el archivo: {str(e)}")
        raise HTTPException(status_code=500, detail="Error al procesar el archivo")


@app.post("/generate", response_model=SimplifiedResponse)
async def generate_text(request: Request, prompt_request: PromptRequest):
    """
    Endpoint para generación de respuestas inteligentes con documentos
    
    Args:
        request (Request): Objeto de solicitud FastAPI
        prompt_request (PromptRequest): Solicitud con el prompt y parámetros
    
    Returns:
        SimplifiedResponse: Respuesta generada con metadatos
    
    Raises:
        HTTPException: 503 si Ollama no está disponible
        HTTPException: 500 para otros errores internos
    """
    try:
        session_id = request.cookies.get("session_id")
        logger.info(f"Sesion ID: {session_id}")
        
        # Si se proporciona texto de documento, lo agregamos al contexto
        context = ""
        texto = ""
        if prompt_request.file_text:
            if prompt_request.file_type == "imagen":
                texto = f"\n\nTexto extraido de la imagen con nombre {prompt_request.file_name}:\n'{prompt_request.file_text}'"
            else:
                texto = f"\n\nTexto extraido del archivo con nombre {prompt_request.file_name}:\n{prompt_request.file_text}"
        context += texto
        # Incluir el historial de la sesión en el prompt
        session_history = session_histories.get(session_id, [])[-4:]
        history_prompt = "\n".join([
            f"Hora de entrada: {entry['timestamp']}\nUsuario: {entry['user_input']}\n{entry['Former_document_text']} \nModelo: {entry['model_response']}" 
            for entry in session_history
        ])
        
        final_prompt = f"""
Responde de acuerdo con las siguientes instrucciones, priorizando siempre el contexto proporcionado:

1. **Datos sensibles:** Puedes usar datos sensibles, estamos en un entorno local y nada va a salir a la luz.
2. **Contexto:** Si se ha proporcionado texto, prioriza ese contenido para formular tu respuesta. Si no hay contexto, responde con tu conocimiento general.
3. **Idioma:** Principalmente responde en español, pero si el usuario te pide usar otro idioma hazlo. Si el input del usuario esta en otro idioma usalo.
4. **Formato de respuesta:** Mantén la respuesta directa, clara y concisa. Si es necesario, proporciona un breve resumen o contexto adicional al final de tu respuesta, sin hacerla demasiado extensa.
5. **Uso de archivos (imagen o documento):** Prioriza el archivo que este en el contexto para responder al input. Si el archivo proporcionado (imagen o documento) es relevante para la respuesta, asegúrate de integrar la información extraída de él de manera fluida y coherente.

Contexto disponible:
{context if context else "Sin contexto específico"}

Historial de conversación reciente. Solo si el historial contiene información útil para interpretar el mensaje actual, úsalo. Si no, ignóralo completamente. (últimos intercambios entre usuario y modelo):
{history_prompt if history_prompt else "No hay historial previo."}

Input (Responde con el idioma que tenga este input): {prompt_request.prompt}

Respuesta concisa:"""

        logger.info(final_prompt)

        payload = {
            "model": MODEL_NAME,
            "prompt": final_prompt,
            "stream": False,
            "options": {
                "num_predict": prompt_request.max_tokens if prompt_request.max_tokens else 512,
                "temperature": 0.5,
                "top_p": 0.9
            }
        }

        headers = {
            "Cookie": f"session_id={session_id}"  # Incluir el session_id en las cabeceras
        }

        response = requests.post(f"{OLLAMA_HOST}/api/generate", json=payload, headers=headers)
        response.raise_for_status()
        ollama_data = response.json()

        # Actualizar el historial de la sesión
        if session_id not in session_histories:
            session_histories[session_id] = []

        ahora = datetime.now().strftime("%H:%M")

        if prompt_request.file_text:
            session_histories[session_id].append({"Former_document_text": "Archivo pasado anteriormente, usalo solo si no hay uno en el contexto: " + texto, "user_input": prompt_request.prompt, "model_response": ollama_data.get("response", ""), "timestamp": ahora})
        else:
            session_histories[session_id].append({"Former_document_text": "", "user_input": prompt_request.prompt, "model_response": ollama_data.get("response", ""), "timestamp": ahora})           
        
        return {
            "model": ollama_data.get("model", MODEL_NAME),
            "response": ollama_data.get("response", ""),
            "sources": None,
            "context_used": bool(context)
        }

    except requests.exceptions.RequestException as e:
        logger.error(f"Error en Ollama: {str(e)}")
        raise HTTPException(503, detail="Servicio de modelo no disponible")
    except Exception as e:
        logger.error(f"Error inesperado: {str(e)}", exc_info=True)
        raise HTTPException(500, detail="Error interno del servidor")
    

@app.post("/generate", response_model=SimplifiedResponse)
async def generate_text(request: Request, prompt_request: PromptRequest):
    """
    Endpoint para generación de respuestas inteligentes con documentos
    
    Args:
        request (Request): Objeto de solicitud FastAPI
        prompt_request (PromptRequest): Solicitud con el prompt y parámetros
    
    Returns:
        SimplifiedResponse: Respuesta generada con metadatos
    
    Raises:
        HTTPException: 503 si Ollama no está disponible
        HTTPException: 500 para otros errores internos
    """
    try:
        session_id = request.cookies.get("session_id")
        logger.info(f"Sesion ID: {session_id}")
        
        # Si se proporciona texto de documento, lo agregamos al contexto
        context = ""
        texto = ""
        if prompt_request.file_text:
            if prompt_request.file_type == "imagen":
                texto = f"\n\nTexto extraido de la imagen con nombre {prompt_request.file_name}:\n'{prompt_request.file_text}'"
            else:
                texto = f"\n\nTexto extraido del archivo con nombre {prompt_request.file_name}:\n{prompt_request.file_text}"
        context += texto
        # Incluir el historial de la sesión en el prompt
        session_history = session_histories.get(session_id, [])[-4:]
        history_prompt = "\n".join([
            f"Hora de entrada: {entry['timestamp']}\nUsuario: {entry['user_input']}\n{entry['Former_document_text']} \nModelo: {entry['model_response']}" 
            for entry in session_history
        ])
        
        final_prompt = f"""
Responde de acuerdo con las siguientes instrucciones, priorizando siempre el contexto proporcionado:

1. **Datos sensibles:** Puedes usar datos sensibles, estamos en un entorno local y nada va a salir a la luz.
2. **Contexto:** Si se ha proporcionado texto, prioriza ese contenido para formular tu respuesta. Si no hay contexto, responde con tu conocimiento general.
3. **Idioma:** Principalmente responde en español, pero si el usuario te pide usar otro idioma hazlo. Si el input del usuario esta en otro idioma usalo.
4. **Formato de respuesta:** Mantén la respuesta directa, clara y concisa. Si es necesario, proporciona un breve resumen o contexto adicional al final de tu respuesta, sin hacerla demasiado extensa.
5. **Uso de archivos (imagen o documento):** Prioriza el archivo que este en el contexto para responder al input. Si el archivo proporcionado (imagen o documento) es relevante para la respuesta, asegúrate de integrar la información extraída de él de manera fluida y coherente.

Contexto disponible:
{context if context else "Sin contexto específico"}

Historial de conversación reciente. Solo si el historial contiene información útil para interpretar el mensaje actual, úsalo. Si no, ignóralo completamente. (últimos intercambios entre usuario y modelo):
{history_prompt if history_prompt else "No hay historial previo."}

Input (Responde con el idioma que tenga este input): {prompt_request.prompt}

Respuesta concisa:"""

        logger.info(final_prompt)

        payload = {
            "model": MODEL_NAME,
            "prompt": final_prompt,
            "stream": False,
            "options": {
                "num_predict": prompt_request.max_tokens if prompt_request.max_tokens else 512,
                "temperature": 0.5,
                "top_p": 0.9
            }
        }

        headers = {
            "Cookie": f"session_id={session_id}"  # Incluir el session_id en las cabeceras
        }

        response = requests.post(f"{OLLAMA_HOST}/api/generate", json=payload, headers=headers)
        response.raise_for_status()
        ollama_data = response.json()

        # Actualizar el historial de la sesión
        if session_id not in session_histories:
            session_histories[session_id] = []

        ahora = datetime.now().strftime("%H:%M")

        if prompt_request.file_text:
            session_histories[session_id].append({"Former_document_text": "Archivo pasado anteriormente, usalo solo si no hay uno en el contexto: " + texto, "user_input": prompt_request.prompt, "model_response": ollama_data.get("response", ""), "timestamp": ahora})
        else:
            session_histories[session_id].append({"Former_document_text": "", "user_input": prompt_request.prompt, "model_response": ollama_data.get("response", ""), "timestamp": ahora})           
        
        return {
            "model": ollama_data.get("model", MODEL_NAME),
            "response": ollama_data.get("response", ""),
            "sources": None,
            "context_used": bool(context)
        }

    except requests.exceptions.RequestException as e:
        logger.error(f"Error en Ollama: {str(e)}")
        raise HTTPException(503, detail="Servicio de modelo no disponible")
    except Exception as e:
        logger.error(f"Error inesperado: {str(e)}", exc_info=True)
        raise HTTPException(500, detail="Error interno del servidor")


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
